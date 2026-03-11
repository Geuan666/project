#!/usr/bin/env python
from __future__ import annotations

import csv
import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/root/autodl-tmp/project")
RESULTS = ROOT / "Automatic-Circuit-Discovery" / "experiments" / "results"
DECK_PATH = Path("/root/autodl-tmp/How does Agentic LLM Choose to Use Tools.pptx")
BACKUP_PATH = Path("/root/autodl-tmp/How does Agentic LLM Choose to Use Tools.backup_before_refresh.pptx")


BG = RGBColor(0xF7, 0xF3, 0xEC)
PAPER = RGBColor(0xFF, 0xFF, 0xFB)
NAVY = RGBColor(0x10, 0x2A, 0x43)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x5B, 0x66, 0x74)
TEAL = RGBColor(0x12, 0x74, 0x75)
TEAL_SOFT = RGBColor(0xDE, 0xF1, 0xEF)
AMBER = RGBColor(0xC2, 0x41, 0x0C)
AMBER_SOFT = RGBColor(0xFC, 0xEA, 0xD9)
SAGE = RGBColor(0x5F, 0x7D, 0x3B)
SAGE_SOFT = RGBColor(0xE5, 0xEF, 0xDB)
SLATE_SOFT = RGBColor(0xE9, 0xEF, 0xF5)
BOR = RGBColor(0xD8, 0xDF, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"


def rgb(fill):
    return fill


def set_background(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color=INK,
    bold=False,
    font_name=BODY_FONT,
    align=PP_ALIGN.LEFT,
    margin=0.08,
    vertical=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return shape


def add_paragraph_box(
    slide,
    left,
    top,
    width,
    height,
    lines,
    font_size=18,
    color=INK,
    font_name=BODY_FONT,
    bullet=False,
    level=0,
    spacing=4,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = level
        p.space_after = Pt(spacing)
        if bullet:
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name
        else:
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name
    return shape


def add_header(slide, section, title, page, subtitle=None):
    set_background(slide)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.22), NAVY, NAVY, 0)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.45), Inches(1.9), Inches(0.38), TEAL, TEAL, 0)
    add_textbox(
        slide,
        Inches(0.68),
        Inches(0.49),
        Inches(1.7),
        Inches(0.28),
        section.upper(),
        14,
        color=WHITE,
        bold=True,
        font_name=BODY_FONT,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(slide, Inches(0.6), Inches(0.92), Inches(10.6), Inches(0.62), title, 26, color=NAVY, bold=True, font_name=TITLE_FONT)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(1.5), Inches(10.8), Inches(0.42), subtitle, 12, color=MUTED)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.95), Inches(12.1), Inches(0.03), TEAL, TEAL, 0)
    add_textbox(
        slide,
        Inches(12.32),
        Inches(7.0),
        Inches(0.42),
        Inches(0.2),
        str(page),
        11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
    add_textbox(
        slide,
        Inches(0.62),
        Inches(7.0),
        Inches(4.4),
        Inches(0.2),
        "Qwen3-1.7B tool-call circuit reproduction",
        10,
        color=MUTED,
    )


def add_card(slide, left, top, width, height, title, body_lines, fill=PAPER, accent=NAVY, body_size=16):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill, BOR, 1.0)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), height, accent, accent, 0)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.14), width - Inches(0.3), Inches(0.3), title, 16, color=accent, bold=True)
    add_paragraph_box(slide, left + Inches(0.18), top + Inches(0.48), width - Inches(0.28), height - Inches(0.58), body_lines, font_size=body_size, color=INK)


def add_metric_card(slide, left, top, width, height, value, label, note="", fill=TEAL, value_color=WHITE):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill, fill, 0)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.12), width - Inches(0.24), Inches(0.46), value, 28, color=value_color, bold=True, font_name=TITLE_FONT)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.62), width - Inches(0.24), Inches(0.28), label, 12, color=WHITE, bold=True)
    if note:
        add_textbox(slide, left + Inches(0.12), top + Inches(0.95), width - Inches(0.24), height - Inches(1.05), note, 10, color=WHITE)


def add_chip(slide, left, top, width, text, fill, color=WHITE):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.38), fill, fill, 0)
    add_textbox(slide, left + Inches(0.06), top + Inches(0.05), width - Inches(0.12), Inches(0.24), text, 11, color=color, bold=True, align=PP_ALIGN.CENTER)


def draw_bar_row(slide, x, y, total_w, label, value, max_value, fill):
    add_textbox(slide, x, y, Inches(2.05), Inches(0.24), label, 12, color=INK, bold=True)
    bar_left = x + Inches(2.15)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, bar_left, y + Inches(0.03), total_w - Inches(2.85), Inches(0.16), SLATE_SOFT, SLATE_SOFT, 0)
    inner_w = (total_w - Inches(2.85)) * (value / max_value if max_value else 0)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, bar_left, y + Inches(0.03), inner_w, Inches(0.16), fill, fill, 0)
    add_textbox(slide, x + total_w - Inches(0.6), y, Inches(0.55), Inches(0.24), f"{value:.3f}", 12, color=fill, bold=True, align=PP_ALIGN.RIGHT)


def load_results():
    with open(ROOT / "datasets" / "clean" / "selection_summary.json", "r", encoding="utf-8") as f:
        selection = json.load(f)
    with open(RESULTS / "toolcall_project_1189_aggregate" / "global_core_summary.json", "r", encoding="utf-8") as f:
        global_core = json.load(f)
    node_roles = []
    with open(RESULTS / "toolcall_project_1189_semantic_roles" / "node_roles.csv", "r", encoding="utf-8") as f:
        node_roles = list(csv.DictReader(f))
    role_groups = []
    with open(RESULTS / "toolcall_project_1189_semantic_roles" / "role_group_summary.csv", "r", encoding="utf-8") as f:
        role_groups = list(csv.DictReader(f))
    with open(RESULTS / "toolcall_project_1189_semantic_roles" / "contrast_token_trace_report.json", "r", encoding="utf-8") as f:
        trace = json.load(f)
    edge_rows = []
    with open(RESULTS / "toolcall_project_1189_semantic_roles" / "path_patch_edge_summary_trimmed.csv", "r", encoding="utf-8") as f:
        edge_rows = list(csv.DictReader(f))
    with open(RESULTS / "toolcall_project_1189_refined_consistent" / "refined_report.json", "r", encoding="utf-8") as f:
        refined = json.load(f)
    consistency_json = RESULTS / "toolcall_project_1189_consistency_eval.json"
    return {
        "selection": selection,
        "global_core": global_core,
        "node_roles": node_roles,
        "role_groups": {row["group"]: row for row in role_groups},
        "trace": trace,
        "edges": edge_rows,
        "refined": refined,
        "consistency_exists": consistency_json.exists(),
    }


def build_deck(data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    selection = data["selection"]
    core = data["global_core"]
    roles = {row["node"]: row for row in data["node_roles"]}
    groups = data["role_groups"]
    trace = data["trace"]
    edges = data["edges"]
    refined = data["refined"]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.26), NAVY, NAVY, 0)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.75), Inches(0.65), Inches(3.95), Inches(5.7), NAVY, NAVY, 0)
    add_textbox(slide, Inches(0.72), Inches(0.82), Inches(7.25), Inches(1.45), "How Does Agentic LLM Choose to Use Tools?", 29, color=NAVY, bold=True, font_name=TITLE_FONT)
    add_textbox(slide, Inches(0.74), Inches(2.2), Inches(7.3), Inches(0.82), "Updated with the full tool-call circuit reproduction on 1,189 clean/corrupt prompt pairs.", 17, color=INK, font_name=BODY_FONT)
    add_paragraph_box(
        slide,
        Inches(0.74),
        Inches(3.05),
        Inches(7.2),
        Inches(1.55),
        [
            "Goal: explain why Qwen3-1.7B emits <tool_call> on clean prompts but not on matched corrupt prompts.",
            "Evidence chain: single-sample ACDC mining -> global core -> semantic roles -> role groups -> contrast tracing -> edge path patching -> refinement.",
        ],
        font_size=16,
        color=MUTED,
    )
    add_textbox(slide, Inches(9.15), Inches(1.02), Inches(3.15), Inches(0.44), "Current headline result", 16, color=WHITE, bold=True, font_name=TITLE_FONT)
    add_paragraph_box(
        slide,
        Inches(9.12),
        Inches(1.55),
        Inches(3.15),
        Inches(2.1),
        [
            "The old Humaneval toy caught the late-stage motif.",
            "The full run shows a broader, more stable dispatch circuit with an earlier MLP scaffold and a stronger writer block.",
        ],
        font_size=15,
        color=WHITE,
    )
    add_paragraph_box(
        slide,
        Inches(9.12),
        Inches(4.05),
        Inches(3.08),
        Inches(1.55),
        [
            "Core motif that survives scaling:",
            "L21/L23/L24 routers + L24H6 tag reader + MLP27 writer",
        ],
        font_size=14,
        color=WHITE,
    )
    add_metric_card(slide, Inches(0.72), Inches(5.78), Inches(2.35), Inches(1.15), "1,189", "paired samples", "clean/corrupt matched by filename", TEAL)
    add_metric_card(slide, Inches(3.22), Inches(5.78), Inches(2.35), Inches(1.15), "10 / 21", "core nodes / edges", "global circuit after aggregation", AMBER)
    add_metric_card(slide, Inches(5.72), Inches(5.78), Inches(2.35), Inches(1.15), "0.966 / 0.968", "full-core suff / nec", "role-group causal validation", SAGE)
    add_textbox(slide, Inches(0.74), Inches(7.0), Inches(4.5), Inches(0.2), "Refreshed on Mar 11, 2026", 10, color=MUTED)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Setup", "Dataset, labels, and why this run is qualitatively stronger", 2)
    add_card(
        slide,
        Inches(0.65),
        Inches(2.25),
        Inches(4.15),
        Inches(2.25),
        "Task setup",
        [
            "Clean: the assistant's first generated token is <tool_call>.",
            "Corrupt: the first token is not <tool_call>.",
            "Pairs are matched by identical filenames across clean/ and corrupt/.",
        ],
        fill=PAPER,
        accent=TEAL,
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(4.72),
        Inches(4.15),
        Inches(1.45),
        "Generalization work",
        [
            "No q1..q164 assumptions, no fixed contrast position, and no single-token difference assumption.",
            "Supports multi-token phrase changes such as write out and directly build.",
        ],
        fill=PAPER,
        accent=AMBER,
        body_size=14,
    )
    add_metric_card(slide, Inches(5.15), Inches(2.28), Inches(2.08), Inches(1.08), "986", "CodeContests", "", NAVY)
    add_metric_card(slide, Inches(7.45), Inches(2.28), Inches(1.72), Inches(1.08), "102", "APPS", "", TEAL)
    add_metric_card(slide, Inches(9.36), Inches(2.28), Inches(1.72), Inches(1.08), "61", "MBPP", "", SAGE)
    add_metric_card(slide, Inches(11.27), Inches(2.28), Inches(1.45), Inches(1.08), "40", "HumanEval", "", AMBER)
    add_metric_card(slide, Inches(5.15), Inches(3.68), Inches(2.28), Inches(1.08), "584", "Java", "", AMBER)
    add_metric_card(slide, Inches(7.66), Inches(3.68), Inches(2.28), Inches(1.08), "402", "C++", "", TEAL)
    add_metric_card(slide, Inches(10.17), Inches(3.68), Inches(2.28), Inches(1.08), "203", "Python", "", SAGE)
    add_card(
        slide,
        Inches(5.15),
        Inches(5.08),
        Inches(7.55),
        Inches(1.1),
        "Trigger vocabulary kept in the final audited set",
        [
            "10 verbs: append, save, insert, paste, write, create, prepare, edit, finalize, prefill",
            "5 phrases: write out, directly build, directly implement, manually build, properly add",
        ],
        fill=PAPER,
        accent=NAVY,
        body_size=13,
    )

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Pipeline", "Full method chain used in the reproduction", 3)
    steps = [
        ("1", "Single-sample ACDC", "Mine local circuits with AP + exact CT on each pair."),
        ("2", "Global aggregation", "Weight and merge circuits into a shared 10-node core."),
        ("3", "Semantic roles", "Measure what each node reads and writes."),
        ("4", "Role-group causality", "Test sufficiency and necessity for functional groups."),
        ("5", "Contrast tracing", "Track where the clean/corrupt difference enters the network."),
        ("6", "Edge path patching", "Quantify mediated edges, not just nodes."),
    ]
    x = Inches(0.67)
    y = Inches(2.28)
    widths = [Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    fills = [TEAL_SOFT, PAPER, SLATE_SOFT, PAPER, AMBER_SOFT, PAPER]
    accents = [TEAL, NAVY, NAVY, NAVY, AMBER, NAVY]
    for idx, (num, title, body) in enumerate(steps):
        w = widths[idx]
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(2.2), fills[idx], BOR, 1.0)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.18), Inches(0.42), Inches(0.34), accents[idx], accents[idx], 0)
        add_textbox(slide, x + Inches(0.18), y + Inches(0.20), Inches(0.36), Inches(0.24), num, 14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.16), y + Inches(0.7), w - Inches(0.32), Inches(0.45), title, 16, color=accents[idx], bold=True)
        add_paragraph_box(slide, x + Inches(0.16), y + Inches(1.15), w - Inches(0.32), Inches(0.82), [body], font_size=13, color=INK)
        x += w + Inches(0.16)
    add_card(
        slide,
        Inches(0.68),
        Inches(5.08),
        Inches(12.0),
        Inches(1.1),
        "What makes this version robust",
        [
            "The tool-call positions are located dynamically, contrast spans can be 1-3 tokens long, and the Input Embed patch is applied to the entire contrast token set rather than a hard-coded single position."
        ],
        fill=PAPER,
        accent=SAGE,
        body_size=14,
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Global Core", "The shared circuit that survives across 1,189 pairs", 4)
    add_metric_card(slide, Inches(0.7), Inches(2.32), Inches(2.0), Inches(1.05), "1,189", "samples used", "", NAVY)
    add_metric_card(slide, Inches(2.88), Inches(2.32), Inches(2.0), Inches(1.05), "10", "core nodes", "", TEAL)
    add_metric_card(slide, Inches(5.06), Inches(2.32), Inches(2.0), Inches(1.05), "21", "core edges", "", AMBER)
    add_card(
        slide,
        Inches(0.72),
        Inches(3.75),
        Inches(12.0),
        Inches(2.2),
        "Backbone flow",
        [
            "Input Embed -> MLP11 -> MLP16 -> {MLP17, L17H8, MLP19} -> {L21H1, L21H12, L23H6, L24H6} -> MLP27 -> Output(<tool_call>)",
            "This keeps the late-stage Humaneval motif but adds a broader early-to-mid MLP scaffold that was invisible in the toy run.",
        ],
        fill=PAPER,
        accent=NAVY,
        body_size=17,
    )
    add_chip(slide, Inches(7.5), Inches(2.42), Inches(2.2), "Late routers survive scaling", TEAL)
    add_chip(slide, Inches(9.88), Inches(2.42), Inches(2.05), "Early scaffold is new", AMBER)
    add_chip(slide, Inches(7.5), Inches(2.92), Inches(4.43), "Global core closes the full decision path", SAGE)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Roles", "Semantic role map for the final core nodes", 5)
    add_card(
        slide,
        Inches(0.68),
        Inches(2.22),
        Inches(3.0),
        Inches(1.65),
        "Tool-tag reader",
        [
            "L24H6",
            "Reads the tool-call tags and both boosts <tool_call> and suppresses the distractor path.",
        ],
        fill=TEAL_SOFT,
        accent=TEAL,
        body_size=15,
    )
    add_card(
        slide,
        Inches(3.98),
        Inches(2.22),
        Inches(3.2),
        Inches(1.65),
        "Format routers",
        [
            "L17H8, L21H1, L21H12, L23H6",
            "Route formatting and dispatch-state information into the final decision block.",
        ],
        fill=SLATE_SOFT,
        accent=NAVY,
        body_size=15,
    )
    add_card(
        slide,
        Inches(7.48),
        Inches(2.22),
        Inches(3.2),
        Inches(1.65),
        "Primary writer MLPs",
        [
            "MLP16, MLP17, MLP19, MLP27",
            "The full run makes the writer stack much more explicit than the old toy circuit did.",
        ],
        fill=SAGE_SOFT,
        accent=SAGE,
        body_size=15,
    )
    add_card(
        slide,
        Inches(10.98),
        Inches(2.22),
        Inches(1.72),
        Inches(1.65),
        "Support",
        [
            "MLP11",
            "Support suppressor",
        ],
        fill=AMBER_SOFT,
        accent=AMBER,
        body_size=13,
    )
    add_card(
        slide,
        Inches(0.68),
        Inches(4.18),
        Inches(12.02),
        Inches(1.92),
        "Node-level highlights",
        [
            f"MLP27 is the strongest writer: target-logit delta median = {float(roles['MLP27']['target_logit_delta_median']):.1f}.",
            f"L21H12 is the strongest router head writer: target-logit delta median = {float(roles['L21H12']['target_logit_delta_median']):.2f}.",
            f"L24H6 stays interpretable and specific: full_ratio_median = {float(roles['L24H6']['full_ratio_median']):.3f}, role = Tool-Tag Reader / Target+Suppressor.",
        ],
        fill=PAPER,
        accent=NAVY,
        body_size=16,
    )

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Causal Validation", "Role groups are causal, not just descriptive", 6)
    full = groups["full_core"]
    writer = groups["primary_writer_mlp"]
    router = groups["format_router"]
    tag = groups["tool_tag_reader"]
    all_heads = groups["all_heads"]
    all_mlps = groups["all_mlps"]
    add_metric_card(slide, Inches(0.72), Inches(2.25), Inches(2.4), Inches(1.28), f"{float(full['suff_median']):.3f}", "full-core sufficiency", "", NAVY)
    add_metric_card(slide, Inches(3.34), Inches(2.25), Inches(2.4), Inches(1.28), f"{float(full['nec_median']):.3f}", "full-core necessity", "", TEAL)
    add_metric_card(slide, Inches(6.0), Inches(2.25), Inches(3.0), Inches(1.28), f"{float(writer['nec_median']):.3f}", "primary-writer MLP necessity", "", SAGE)
    add_metric_card(slide, Inches(9.22), Inches(2.25), Inches(3.0), Inches(1.28), f"{float(router['drop_full_nec_median']):.3f}", "drop if routers are removed", "", AMBER)
    add_card(
        slide,
        Inches(0.72),
        Inches(3.88),
        Inches(5.95),
        Inches(2.22),
        "Interpretation",
        [
            "The full core nearly saturates the clean/corrupt gap.",
            "Primary writer MLPs are now strong enough to matter on their own.",
            "Format routers are still the main control backbone: remove them and necessity drops sharply.",
        ],
        fill=PAPER,
        accent=NAVY,
        body_size=16,
    )
    add_card(
        slide,
        Inches(6.95),
        Inches(3.88),
        Inches(5.78),
        Inches(2.22),
        "Heads vs MLPs",
        [
            f"all_heads: suff {float(all_heads['suff_median']):.3f}, nec {float(all_heads['nec_median']):.3f}",
            f"all_mlps: suff {float(all_mlps['suff_median']):.3f}, nec {float(all_mlps['nec_median']):.3f}",
            "Unlike the old Humaneval toy, this full run is no longer head-only. The decision path is now head + MLP co-driven.",
        ],
        fill=PAPER,
        accent=TEAL,
        body_size=16,
    )
    add_chip(slide, Inches(8.55), Inches(2.0), Inches(1.75), "writer-heavy", SAGE)
    add_chip(slide, Inches(10.5), Inches(2.0), Inches(1.95), "router-critical", AMBER)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Tracing", "Contrast tracing shows where the decision signal first appears", 7)
    hist = trace["contrast_position_hist"]
    span_hist = trace["contrast_span_length_hist"]
    add_card(
        slide,
        Inches(0.72),
        Inches(2.25),
        Inches(3.45),
        Inches(1.58),
        "Dynamic contrast positions",
        [
            f"134: {hist.get('134', 0)}  |  135: {hist.get('135', 0)}",
            f"136: {hist.get('136', 0)}  |  137: {hist.get('137', 0)}",
            "The new data does not collapse to one hard-coded position.",
        ],
        fill=PAPER,
        accent=AMBER,
        body_size=15,
    )
    add_card(
        slide,
        Inches(0.72),
        Inches(4.02),
        Inches(3.45),
        Inches(1.58),
        "Dynamic span lengths",
        [
            f"1 token: {span_hist.get('1', 0)}",
            f"2 tokens: {span_hist.get('2', 0)}",
            f"3 tokens: {span_hist.get('3', 0)}",
        ],
        fill=PAPER,
        accent=TEAL,
        body_size=15,
    )
    add_card(
        slide,
        Inches(4.45),
        Inches(2.25),
        Inches(8.25),
        Inches(3.35),
        "Early-layer recovery profile",
        [
            "Median contrast recovery is strongest at the earliest layers, then decays quickly.",
            "tool_call-tag and prefix-only patching stay near zero at the same layers.",
        ],
        fill=PAPER,
        accent=NAVY,
        body_size=16,
    )
    first10 = trace["contrast_recovery_median"][:10]
    max_first10 = max(first10)
    for idx, val in enumerate(first10):
        draw_bar_row(slide, Inches(4.85), Inches(3.08 + idx * 0.22), Inches(7.15), f"L{idx}", float(val), max_first10, AMBER)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Path Patching", "Top mediated edges in the final global circuit", 8)
    top_edges = [
        ("MLP27 -> Output", float(next(r for r in edges if r["edge"] == "MLP27->Residual Output: <tool_call>")["edge_ratio_median"])),
        ("L24H6 -> Output", float(next(r for r in edges if r["edge"] == "L24H6->Residual Output: <tool_call>")["edge_ratio_median"])),
        ("L23H6 -> Output", float(next(r for r in edges if r["edge"] == "L23H6->Residual Output: <tool_call>")["edge_ratio_median"])),
        ("L21H12 -> MLP27", float(next(r for r in edges if r["edge"] == "L21H12->MLP27")["edge_ratio_median"])),
        ("L24H6 -> MLP27", float(next(r for r in edges if r["edge"] == "L24H6->MLP27")["edge_ratio_median"])),
        ("L21H1 -> MLP27", float(next(r for r in edges if r["edge"] == "L21H1->MLP27")["edge_ratio_median"])),
        ("MLP16 -> MLP17", float(next(r for r in edges if r["edge"] == "MLP16->MLP17")["edge_ratio_median"])),
        ("MLP16 -> MLP19", float(next(r for r in edges if r["edge"] == "MLP16->MLP19")["edge_ratio_median"])),
    ]
    max_edge = max(v for _, v in top_edges)
    for idx, (label, value) in enumerate(top_edges):
        draw_bar_row(slide, Inches(0.84), Inches(2.33 + idx * 0.42), Inches(11.2), label, value, max_edge, TEAL if idx < 3 else AMBER)
    add_card(
        slide,
        Inches(9.45),
        Inches(5.55),
        Inches(3.2),
        Inches(0.75),
        "Main reading",
        [
            "The final write-out is dominated by MLP27, but the L21/L23/L24 route is directly measurable upstream."
        ],
        fill=PAPER,
        accent=NAVY,
        body_size=12,
    )

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Consistency", "Refinement makes the circuit stable without destroying performance", 9)
    add_metric_card(slide, Inches(0.75), Inches(2.32), Inches(2.65), Inches(1.25), f"{refined['consistency']['orig_pairwise_jaccard_mean']:.3f}", "orig Jaccard mean", "", AMBER)
    add_metric_card(slide, Inches(3.62), Inches(2.32), Inches(2.65), Inches(1.25), f"{refined['consistency']['refined_pairwise_jaccard_mean']:.3f}", "refined Jaccard mean", "", TEAL)
    add_metric_card(slide, Inches(6.5), Inches(2.32), Inches(2.65), Inches(1.25), f"{refined['consistency']['orig_pairwise_jaccard_median']:.3f}", "orig Jaccard median", "", AMBER)
    add_metric_card(slide, Inches(9.38), Inches(2.32), Inches(2.65), Inches(1.25), f"{refined['consistency']['refined_pairwise_jaccard_median']:.1f}", "refined Jaccard median", "", TEAL)
    add_card(
        slide,
        Inches(0.75),
        Inches(4.0),
        Inches(5.9),
        Inches(1.9),
        "Performance before / after refinement",
        [
            f"orig: suff {refined['orig_metrics']['suff_median']:.3f}, nec {refined['orig_metrics']['nec_median']:.3f}",
            f"refined: suff {refined['refined_metrics']['suff_median']:.3f}, nec {refined['refined_metrics']['nec_median']:.3f}",
            "The stable backbone is not a cosmetic cleanup. It keeps almost all of the circuit's causal power.",
        ],
        fill=PAPER,
        accent=NAVY,
        body_size=16,
    )
    add_card(
        slide,
        Inches(6.95),
        Inches(4.0),
        Inches(5.48),
        Inches(1.9),
        "Refinement output",
        [
            "Stable 10-node backbone:",
            ", ".join(refined["backbone_nodes"]),
            "At most 4 local nodes are retained on top of the shared backbone.",
        ],
        fill=PAPER,
        accent=SAGE,
        body_size=15,
    )
    if not data["consistency_exists"]:
        add_textbox(slide, Inches(0.82), Inches(6.26), Inches(8.2), Inches(0.25), "Note: refined consistency metrics are complete; the final convenience JSON export is still running.", 10, color=MUTED)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Takeaways", "What the full run adds beyond the old Humaneval toy", 10)
    add_card(
        slide,
        Inches(0.72),
        Inches(2.22),
        Inches(3.75),
        Inches(3.8),
        "1. Scale and data realism",
        [
            "Toy: 139 valid samples after filtering, mostly a template-stable Humaneval-style regime.",
            "Now: 1,189 audited pairs across CodeContests, APPS, MBPP, and HumanEval with verb and phrase edits.",
        ],
        fill=PAPER,
        accent=TEAL,
        body_size=16,
    )
    add_card(
        slide,
        Inches(4.8),
        Inches(2.22),
        Inches(3.75),
        Inches(3.8),
        "2. Circuit structure",
        [
            "Toy already found the late-stage motif: L21/L23/L24 + MLP27.",
            "The full run adds an earlier scaffold centered on MLP16/17/19 instead of a single L20H5 bottleneck.",
        ],
        fill=PAPER,
        accent=AMBER,
        body_size=16,
    )
    add_card(
        slide,
        Inches(8.88),
        Inches(2.22),
        Inches(3.75),
        Inches(3.8),
        "3. Causal story",
        [
            "Toy looked head-heavy.",
            "The full run is head + MLP co-driven, with a stronger and more explicit writer stack.",
            "This is closer to a real tool-call dispatch circuit than a prompt-template artifact.",
        ],
        fill=PAPER,
        accent=SAGE,
        body_size=16,
    )
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(6.28), Inches(11.92), Inches(0.56), NAVY, NAVY, 0)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(6.42),
        Inches(11.5),
        Inches(0.2),
        "Bottom line: we now have a cross-sample, semantically labeled, edge-validated tool-call circuit with strong causal sufficiency and necessity.",
        14,
        color=WHITE,
        bold=True,
        font_name=BODY_FONT,
        align=PP_ALIGN.CENTER,
    )

    return prs


def verify_deck(path: Path):
    prs = Presentation(path)
    picture_count = 0
    titles = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        titles.append(text)
                        break
                break
            if getattr(shape, "shape_type", None) == 13:
                picture_count += 1
    media_count = 0
    import zipfile

    with zipfile.ZipFile(path) as zf:
        media_count = sum(1 for name in zf.namelist() if name.startswith("ppt/media/"))
    return {
        "slides": len(prs.slides),
        "picture_count": picture_count,
        "media_count": media_count,
        "first_texts": titles[:10],
    }


def main():
    data = load_results()
    if DECK_PATH.exists() and not BACKUP_PATH.exists():
        shutil.copy2(DECK_PATH, BACKUP_PATH)
    prs = build_deck(data)
    prs.save(DECK_PATH)
    report = verify_deck(DECK_PATH)
    print(json.dumps(report, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
